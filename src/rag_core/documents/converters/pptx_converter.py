"""PPTX converter with hybrid text extraction + OCR fallback.

Extracts text from slide shapes, tables, and speaker notes.
"""

from __future__ import annotations

import asyncio
import io
import logging
from dataclasses import dataclass
from typing import Any, Dict, List

from .base import (
    ConversionResult,
    HybridConverter,
    render_markdown_table,
    score_text_quality,
)
from .converter_keys import PPTX_CONVERTER_KEY

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class _ShapeGeometry:
    left: int
    top: int
    width: int
    height: int


@dataclass(frozen=True)
class _SlideTextAnchor:
    section_start: int
    text: str
    geometry: _ShapeGeometry
    is_title: bool


def _is_generic_figure_description(value: str) -> bool:
    normalized = " ".join(value.lower().split())
    if normalized.startswith("picture "):
        suffix = normalized.removeprefix("picture ").strip()
        return suffix.isdigit()
    if normalized.startswith("image "):
        suffix = normalized.removeprefix("image ").strip()
        return suffix.isdigit()
    return False


def _extract_shape_text(shape: Any) -> List[str]:
    """Extract text lines from a PPTX shape (text frames and tables)."""
    lines: List[str] = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)

    if shape.has_table:
        rows: List[List[str]] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            lines.append(render_markdown_table(rows))

    return lines


def _extract_slide_figure_items(
    slide: Any,
    slide_index: int,
) -> tuple[List[Dict[str, Any]], dict[str, _ShapeGeometry]]:
    """Extract figure metadata from image-like slide shapes."""
    figures: List[Dict[str, Any]] = []
    geometries: dict[str, _ShapeGeometry] = {}
    picture_shape_type = None
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        picture_shape_type = MSO_SHAPE_TYPE.PICTURE
    except Exception:
        # python-pptx may be unavailable in some environments; keep extraction going.
        pass

    figure_number = 0
    for shape in slide.shapes:
        is_picture = False
        try:
            if (
                picture_shape_type is not None
                and shape.shape_type == picture_shape_type
            ):
                is_picture = True
        except Exception:
            # Some shapes expose shape_type inconsistently; fall back to image detection.
            pass
        if not is_picture and getattr(shape, "image", None) is not None:
            is_picture = True
        if not is_picture:
            continue

        figure_number += 1
        figure_id = "fig:slide:%d:%d" % (slide_index + 1, figure_number)
        label = "Slide %d Figure %d" % (slide_index + 1, figure_number)
        description = ""
        try:
            alt_text = str(getattr(shape, "name", "") or "").strip()
            if alt_text and not _is_generic_figure_description(alt_text):
                description = alt_text
        except Exception:
            # Alt text is optional; missing metadata should not block figure extraction.
            pass

        figures.append(
            {
                "figure_id": figure_id,
                "page_index": slide_index,
                "label": label,
                "description": description,
                "metadata": {
                    "source": "pptx:picture_shape",
                    "slide_number": slide_index + 1,
                },
            }
        )
        geometries[figure_id] = _shape_geometry(shape)

    return figures, geometries


def _attach_single_pptx_figure_locators(
    figure_items: List[Dict[str, Any]],
    *,
    content: str,
    slide_sections: List[str],
    text_anchors: dict[int, list[_SlideTextAnchor]],
    figure_geometries: dict[str, _ShapeGeometry],
) -> None:
    figures_by_slide: dict[int, list[Dict[str, Any]]] = {}
    for item in figure_items:
        metadata = item.get("metadata")
        if not isinstance(metadata, dict):
            continue
        slide_number = metadata.get("slide_number")
        if isinstance(slide_number, int):
            figures_by_slide.setdefault(slide_number, []).append(item)

    cursor = 0
    separator_length = len("\n\n---\n\n")
    for slide_number, section in enumerate(slide_sections, start=1):
        figures = figures_by_slide.get(slide_number, [])
        section_end = cursor + len(section)
        if (
            len(figures) == 1
            and not str(figures[0].get("description", "") or "").strip()
        ):
            figure_id = str(figures[0].get("figure_id", "") or "")
            figure_geometry = figure_geometries.get(figure_id)
            candidate = (
                _nearest_text_anchor(
                    figure_geometry,
                    text_anchors.get(slide_number, []),
                )
                if figure_geometry is not None
                else None
            )
            if candidate is not None:
                anchor_start = cursor + candidate.section_start
                metadata = dict(figures[0].get("metadata") or {})
                anchor_point = anchor_start + max(0, len(candidate.text) // 2)
                metadata["text_anchor_start_char"] = anchor_point
                metadata["text_anchor_end_char"] = anchor_point
                figures[0]["metadata"] = metadata
        cursor = section_end + separator_length


def _shape_geometry(shape: Any) -> _ShapeGeometry:
    return _ShapeGeometry(
        left=int(getattr(shape, "left", 0) or 0),
        top=int(getattr(shape, "top", 0) or 0),
        width=int(getattr(shape, "width", 0) or 0),
        height=int(getattr(shape, "height", 0) or 0),
    )


def _nearest_text_anchor(
    figure: _ShapeGeometry,
    anchors: list[_SlideTextAnchor],
) -> _SlideTextAnchor | None:
    candidates = [anchor for anchor in anchors if not anchor.is_title]
    if not candidates:
        return None
    ranked = sorted(
        (
            (_rectangle_gap_squared(figure, anchor.geometry), anchor)
            for anchor in candidates
        ),
        key=lambda item: item[0],
    )
    if len(ranked) > 1 and ranked[0][0] == ranked[1][0]:
        return None
    return ranked[0][1]


def _rectangle_gap_squared(
    first: _ShapeGeometry,
    second: _ShapeGeometry,
) -> int:
    horizontal = max(
        first.left - (second.left + second.width),
        second.left - (first.left + first.width),
        0,
    )
    vertical = max(
        first.top - (second.top + second.height),
        second.top - (first.top + first.height),
        0,
    )
    return horizontal * horizontal + vertical * vertical


class PptxConverter(HybridConverter):
    """Converts PPTX files to markdown with slide structure and speaker notes."""

    format_name = PPTX_CONVERTER_KEY

    async def convert(
        self,
        file_bytes: bytes,
        filename: str,
        mime_type: str,
    ) -> ConversionResult:
        return await self._convert_office(
            file_bytes, filename, mime_type, parse_label="PPTX"
        )

    async def _try_extract(
        self,
        file_bytes: bytes,
        filename: str,
        mime_type: str,
    ) -> ConversionResult:
        """Extract text from PPTX using python-pptx."""
        from pptx import Presentation

        def _extract() -> ConversionResult:
            try:
                prs = Presentation(io.BytesIO(file_bytes))
            except Exception as exc:
                logger.warning(
                    "Failed to open Office document: format=%s error_type=%s",
                    self.format_name,
                    type(exc).__name__,
                )
                raise ValueError("PPTX parse failed (%s)" % type(exc).__name__) from exc

            slide_sections: List[str] = []
            figure_items: List[Dict[str, Any]] = []
            extracted_text_parts: List[str] = []
            slide_text_anchors: dict[int, list[_SlideTextAnchor]] = {}
            figure_geometries: dict[str, _ShapeGeometry] = {}

            for i, slide in enumerate(prs.slides):
                parts: List[str] = ["## Slide %d" % (i + 1)]
                try:
                    title_shape = slide.shapes.title
                except Exception:
                    # Some slide masters omit a readable title placeholder; keep extraction going without it.
                    title_shape = None
                slide_title = (
                    str(getattr(title_shape, "text", "") or "").strip()
                    if title_shape is not None
                    else ""
                )
                if slide_title:
                    parts.append("### %s" % slide_title)
                    extracted_text_parts.append(slide_title)

                anchors: list[_SlideTextAnchor] = []
                section_length = sum(len(part) for part in parts) + 2 * (len(parts) - 1)
                for shape in slide.shapes:
                    shape_lines = _extract_shape_text(shape)
                    geometry = _shape_geometry(shape)
                    is_title = title_shape is not None and getattr(
                        shape,
                        "element",
                        None,
                    ) is getattr(title_shape, "element", None)
                    shape_start = section_length + (2 if shape_lines and parts else 0)
                    for line in shape_lines:
                        if parts:
                            section_length += 2
                        parts.append(line)
                        section_length += len(line)
                    if shape_lines:
                        anchors.append(
                            _SlideTextAnchor(
                                section_start=shape_start,
                                text="\n\n".join(shape_lines),
                                geometry=geometry,
                                is_title=is_title,
                            )
                        )
                    extracted_text_parts.extend(shape_lines)
                slide_text_anchors[i + 1] = anchors

                slide_figures, slide_figure_geometries = _extract_slide_figure_items(
                    slide, i
                )
                figure_geometries.update(slide_figure_geometries)
                if slide_figures:
                    for item in slide_figures:
                        description = str(item.get("description", "") or "").strip()
                        if description:
                            parts.append(description)
                            extracted_text_parts.append(description)
                    figure_items.extend(slide_figures)

                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append("\n> **Notes:** %s" % notes)
                        extracted_text_parts.append(notes)

                slide_sections.append("\n\n".join(parts))

            text_content = "\n\n".join(extracted_text_parts).strip()
            content = "\n\n---\n\n".join(slide_sections) if text_content else ""
            _attach_single_pptx_figure_locators(
                figure_items,
                content=content,
                slide_sections=slide_sections,
                text_anchors=slide_text_anchors,
                figure_geometries=figure_geometries,
            )
            quality = score_text_quality(content)

            metadata: Dict[str, Any] = {
                "parser": "local:python-pptx",
                "slide_count": len(prs.slides),
                "needs_ocr": bool(figure_items and not text_content),
                "text_char_count": len(text_content),
            }
            if figure_items:
                metadata["figure_items"] = figure_items
                metadata["figure_count"] = len(figure_items)

            return ConversionResult(
                content=content,
                metadata=metadata,
                quality=quality,
            )

        return await asyncio.to_thread(_extract)
