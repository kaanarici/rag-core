from __future__ import annotations

import asyncio
import hashlib
import json
import os
import re
import tracemalloc
import zipfile
from collections import Counter
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from statistics import mean, median
from time import perf_counter
from typing import cast

import pytest

from rag_core._engine.core_lifecycle import resolve_document_id
from rag_core._engine.core_archive_ingest import ingest_zip_archive_with_core
from rag_core._engine.core_prepare import prepare_document_bytes
from rag_core.config import ChunkingConfig
from rag_core.core_models import IngestedDocument, PreparedDocument
from rag_core.documents.local_parse import parse_file_bytes
from rag_core.documents.ocr import OcrRequest, OcrResult
from rag_core.ingest.sources.local import read_local_source
from rag_core.search.policy import DEFAULT_POLICY

pytestmark = [pytest.mark.eval]

_REAL_FIXTURES = (
    Path(__file__).resolve().parents[1]
    / "fixtures"
    / "real_documents"
    / "apache_tika"
)
_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
    b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
    b"\x00\x00\x0cIDATx\x9cc```\x00\x00\x00\x04\x00\x01"
    b"\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
)
_REFERENCE_TOKEN_UNIT_RE = re.compile(r"[A-Za-z0-9_]+|[^\s]", re.UNICODE)


@dataclass(frozen=True)
class _Case:
    name: str
    filename: str
    mime_type: str
    payload: bytes
    gold: str
    locator_key: str | None = None
    locator_value: object | None = None


class _ImageOcr:
    async def extract_markdown(self, request: OcrRequest) -> OcrResult:
        assert request.mime_type == "image/png"
        return OcrResult(
            markdown="# Image record\n\nImage fact beacon-771 survives OCR.",
            provider_name="fixture-ocr",
            model_name="fixture-v1",
            pages_processed=[0],
        )


def _reference_token_units(text: str) -> int:
    count = 0
    for match in _REFERENCE_TOKEN_UNIT_RE.finditer(text):
        unit = match.group(0)
        if unit.isascii() and (unit[0].isalnum() or unit[0] == "_"):
            count += max(1, (len(unit) + 3) // 4)
        else:
            count += len(unit)
    return count


async def _prepare(case: _Case, *, ocr_provider: object | None = None) -> PreparedDocument:
    return await prepare_document_bytes(
        file_bytes=case.payload,
        filename=case.filename,
        mime_type=case.mime_type,
        path=None,
        ocr_provider=ocr_provider,  # type: ignore[arg-type]
    )


async def _prepare_measured(
    case: _Case,
    *,
    ocr_provider: object | None = None,
) -> tuple[PreparedDocument, float]:
    prepared: PreparedDocument | None = None
    latencies: list[float] = []
    for _ in range(3):
        started = perf_counter()
        current = await _prepare(case, ocr_provider=ocr_provider)
        latencies.append((perf_counter() - started) * 1000)
        prepared = prepared or current
    assert prepared is not None
    return prepared, median(latencies)


def _text_cases() -> tuple[_Case, ...]:
    return (
        _Case(
            "html",
            "status.html",
            "text/html",
            (
                b"<html><body><nav>ignore</nav><main><h1>Operations</h1>"
                b"<p>HTML fact heliotrope-113 is retained.</p></main></body></html>"
            ),
            "HTML fact heliotrope-113 is retained.",
            "section_title",
            "Operations",
        ),
        _Case(
            "csv",
            "status.csv",
            "text/csv",
            b"item,status\ncsv-fact-227,ready\nnear-match,waiting\n",
            "csv-fact-227",
        ),
        _Case(
            "tsv",
            "status.tsv",
            "text/tab-separated-values",
            b"item\tstatus\ntsv-fact-331\tready\nnear-match\twaiting\n",
            "tsv-fact-331",
        ),
        _Case(
            "jsonl",
            "status.jsonl",
            "application/x-ndjson",
            (
                b'{"item":"jsonl-fact-449","status":"ready"}\n'
                b'{"item":"near-match","status":"waiting"}\n'
            ),
            "jsonl-fact-449",
        ),
        _Case(
            "code",
            "status.py",
            "text/x-python",
            (
                b"def status_marker() -> str:\n"
                b'    return "code-fact-557 remains addressable"\n'
            ),
            "code-fact-557 remains addressable",
            "line_start",
            1,
        ),
    )


def _real_cases() -> tuple[_Case, ...]:
    return (
        _Case(
            "pdf-real",
            "testPDF.pdf",
            "application/pdf",
            (_REAL_FIXTURES / "testPDF.pdf").read_bytes(),
            "Content Analysis Toolkit",
            "page_number",
            1,
        ),
        _Case(
            "docx-real",
            "testWORD.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            (_REAL_FIXTURES / "testWORD.docx").read_bytes(),
            "The table has things in it",
        ),
        _Case(
            "pptx-real",
            "testPPT.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            (_REAL_FIXTURES / "testPPT.pptx").read_bytes(),
            "Watershed",
            "slide_number",
            3,
        ),
        _Case(
            "xlsx-real",
            "testEXCEL.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            (_REAL_FIXTURES / "testEXCEL.xlsx").read_bytes(),
            "Numbers and their Squares",
            "sheet_name",
            "Feuil1",
        ),
    )


def _mixed_docx_case() -> _Case:
    from docx import Document

    document = Document()
    document.add_heading("Primary architecture", level=1)
    document.add_paragraph("Primary path fact cobalt-601 belongs with the next figure.")
    first = document.add_picture(BytesIO(_PNG))
    first._inline.docPr.attrib["descr"] = "Primary architecture diagram"
    document.add_heading("Recovery architecture", level=1)
    document.add_paragraph("Recovery path fact cobalt-602 belongs with the next figure.")
    second = document.add_picture(BytesIO(_PNG))
    second._inline.docPr.attrib["descr"] = "Recovery architecture diagram"
    buffer = BytesIO()
    document.save(buffer)
    return _Case(
        "docx-mixed",
        "mixed.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        buffer.getvalue(),
        "Recovery path fact cobalt-602 belongs with the next figure.",
    )


def _mixed_pptx_case() -> _Case:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for marker in (
        "Primary slide fact indigo-641 belongs with its figure.",
        "Recovery slide fact indigo-642 belongs with its figure.",
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(0.6),
            Inches(5.5),
            Inches(0.8),
        )
        textbox.text_frame.text = marker
        slide.shapes.add_picture(
            BytesIO(_PNG),
            Inches(0.8),
            Inches(1.7),
            width=Inches(1.0),
            height=Inches(1.0),
        )
    buffer = BytesIO()
    presentation.save(buffer)
    return _Case(
        "pptx-mixed",
        "mixed.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        buffer.getvalue(),
        "Recovery slide fact indigo-642 belongs with its figure.",
        "slide_number",
        2,
    )


def _large_xlsx_case() -> _Case:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Inventory"
    sheet.append(["record", "status", "detail"])
    for index in range(1, 301):
        detail = (
            "quartz-fact-733 retained in the final spreadsheet window"
            if index == 299
            else f"ordinary inventory row {index:03d}"
        )
        sheet.append([f"record-{index:03d}", "ready", detail])
    buffer = BytesIO()
    workbook.save(buffer)
    return _Case(
        "xlsx-large",
        "inventory.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        buffer.getvalue(),
        "quartz-fact-733 retained in the final spreadsheet window",
    )


def _mixed_pdf_bytes() -> bytes:
    import fitz

    document = fitz.open()
    text_page = document.new_page()
    text_page.insert_text(
        (72, 72),
        "Readable PDF fact amber-811 remains on the text page. " * 3,
    )
    image_page = document.new_page()
    image_page.insert_image(fitz.Rect(72, 72, 172, 172), stream=_PNG)
    payload = document.tobytes()
    document.close()
    return bytes(payload)


def _layout_pdf_case() -> _Case:
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text(
        (72, 80),
        "LEFT ONE fact violet-921 is the first left block.",
    )
    page.insert_text(
        (320, 95),
        "RIGHT ONE fact violet-922 is the first right block.",
    )
    page.insert_text(
        (72, 150),
        "LEFT TWO continues the left reading order.",
    )
    page.insert_text(
        (320, 165),
        "RIGHT TWO continues the right reading order.",
    )
    payload = document.tobytes()
    document.close()
    return _Case(
        "pdf-layout",
        "columns.pdf",
        "application/pdf",
        payload,
        "LEFT ONE fact violet-921",
        "page_number",
        1,
    )


async def _archive_metrics(tmp_path: Path) -> dict[str, float | int]:
    member_count = 1_000
    access_path = tmp_path / "access.zip"
    access_payload = b"# Archive access record\n\nbounded member payload"
    with zipfile.ZipFile(access_path, "w", compression=zipfile.ZIP_STORED) as archive:
        for index in range(member_count):
            archive.writestr(f"docs/item-{index:04d}.md", access_payload)

    class AccessCore:
        count = 0

        async def add_bytes(self, **kwargs: object) -> IngestedDocument:
            self.count += 1
            file_bytes = cast(bytes, kwargs["file_bytes"])
            return IngestedDocument(
                document_id=f"archive-eval-{self.count}",
                namespace=cast(str, kwargs["namespace"]),
                collection=cast(str, kwargs["collection"]),
                chunk_count=1,
                filename=cast(str, kwargs["filename"]),
                mime_type=cast(str, kwargs["mime_type"]),
                document_key=cast(str | None, kwargs["document_key"]),
                content_sha256=hashlib.sha256(file_bytes).hexdigest(),
            )

    started = perf_counter()
    access_result = await ingest_zip_archive_with_core(
        core=AccessCore(),
        archive_path=access_path,
        namespace="eval",
        collection="docs",
        max_concurrency=4,
    )
    elapsed_ms = (perf_counter() - started) * 1000

    memory_path = tmp_path / "memory.zip"
    memory_payload = bytes(range(256)) * 2_048
    with zipfile.ZipFile(memory_path, "w", compression=zipfile.ZIP_STORED) as archive:
        for index in range(8):
            archive.writestr(f"docs/item-{index:02d}.md", memory_payload)

    tracemalloc.start()
    memory_result = await ingest_zip_archive_with_core(
        core=AccessCore(),
        archive_path=memory_path,
        namespace="eval",
        collection="docs",
        max_concurrency=4,
    )
    _, peak = tracemalloc.get_traced_memory()
    tracemalloc.stop()
    return {
        "archive_item_count": access_result.succeeded_count,
        "archive_access_latency_ms": round(elapsed_ms, 3),
        "archive_memory_item_count": memory_result.succeeded_count,
        "archive_memory_peak_bytes": peak,
    }


async def collect_document_preparation_metrics(
    tmp_path: Path,
) -> dict[str, float | int | bool]:
    cases = (
        *_real_cases(),
        *_text_cases(),
        _mixed_docx_case(),
        _mixed_pptx_case(),
        _large_xlsx_case(),
    )
    prepared: dict[str, PreparedDocument] = {}
    latencies: list[float] = []
    for case in cases:
        prepared[case.name], latency = await _prepare_measured(case)
        latencies.append(latency)

    layout_case = _layout_pdf_case()
    prepared[layout_case.name], latency = await _prepare_measured(layout_case)
    latencies.append(latency)
    cases = (*cases, layout_case)

    image_case = _Case(
        "image",
        "scan.png",
        "image/png",
        _PNG,
        "Image fact beacon-771 survives OCR.",
    )
    image_markdown, image_parse_metadata = await parse_file_bytes(
        file_bytes=image_case.payload,
        filename=image_case.filename,
        mime_type=image_case.mime_type,
    )
    assert image_markdown == ""
    prepared[image_case.name], latency = await _prepare_measured(
        image_case,
        ocr_provider=_ImageOcr(),
    )
    latencies.append(latency)
    cases = (*cases, image_case)

    mixed_pdf = _mixed_pdf_bytes()
    mixed_markdown, mixed_metadata = await parse_file_bytes(
        file_bytes=mixed_pdf,
        filename="mixed.pdf",
        mime_type="application/pdf",
    )

    dense_text = "# Dense tokens\n\n" + ("界" * 350) + ("+-=" * 50)
    dense = await prepare_document_bytes(
        file_bytes=dense_text.encode(),
        filename="dense.md",
        mime_type="text/markdown",
        path=None,
        ocr_provider=None,
        chunking_config=ChunkingConfig(max_chars=400, overlap=0),
    )
    reference_token_counts = [
        _reference_token_units(chunk.text) for chunk in dense.chunks
    ]
    token_budget = 100

    exact_retained = sum(case.gold in prepared[case.name].markdown for case in cases)
    locator_expected = [case for case in cases if case.locator_key is not None]
    locator_survived = sum(
        any(
            chunk.metadata.get(cast(str, case.locator_key)) == case.locator_value
            and case.gold in chunk.text
            for chunk in prepared[case.name].chunks
        )
        for case in locator_expected
    )
    reliable_chunks = [
        chunk
        for document in prepared.values()
        for chunk in document.chunks
        if chunk.start_char is not None and chunk.end_char is not None
    ]
    reliable_spans = sum(
        chunk.metadata.get("offset_reconstruction") != "unreliable"
        and document.markdown[chunk.start_char : chunk.end_char] == chunk.text
        for document in prepared.values()
        for chunk in document.chunks
        if chunk.start_char is not None and chunk.end_char is not None
    )
    figure_assignments = Counter(
        str(figure_id)
        for name in ("docx-mixed", "pptx-mixed")
        for chunk in prepared[name].chunks
        if (figure_id := chunk.metadata.get("figure_id"))
    )
    large_gold_chunk = next(
        chunk
        for chunk in prepared["xlsx-large"].chunks
        if "quartz-fact-733" in chunk.text
    )
    large_xlsx_chunks = prepared["xlsx-large"].chunks
    large_xlsx_total_tokens = sum(chunk.token_count for chunk in large_xlsx_chunks)
    large_xlsx_repeated_context_tokens = sum(
        value
        for chunk in large_xlsx_chunks[1:]
        if isinstance(
            value := chunk.metadata.get("xlsx_context_token_count"),
            int,
        )
        and not isinstance(value, bool)
    )

    folder = tmp_path / "folder"
    folder.mkdir()
    (folder / "note.md").write_text("folder fact", encoding="utf-8")
    (folder / "scan.png").write_bytes(_PNG)
    folder_plan = read_local_source(folder)

    duplicate = b"identity fact identical bytes"
    near_duplicate = duplicate + b"."
    same_key_id = resolve_document_id(
        namespace="eval",
        collection="docs",
        document_key="shared-key",
        document_id=None,
        policy=DEFAULT_POLICY,
    )
    other_key_id = resolve_document_id(
        namespace="eval",
        collection="docs",
        document_key="other-key",
        document_id=None,
        policy=DEFAULT_POLICY,
    )

    metrics: dict[str, float | int | bool] = {
        "case_count": len(cases),
        "exact_fact_retention": exact_retained / len(cases),
        "locator_survival": locator_survived / len(locator_expected),
        "reliable_span_rate": reliable_spans / len(reliable_chunks),
        "chunk_count": sum(
            len(prepared[case.name].chunks) for case in cases
        ),
        "duplicate_chunk_count": sum(
            len(document.chunks)
            - len({hashlib.sha256(chunk.text.encode()).digest() for chunk in document.chunks})
            for document in prepared.values()
        ),
        "prepare_latency_mean_ms": round(mean(latencies), 3),
        "prepare_latency_max_ms": round(max(latencies), 3),
        "page_specific_ocr_indices": mixed_metadata.get("ocr_page_indices") == [1],
        "mixed_pdf_text_retained": "amber-811" in mixed_markdown,
        "image_requires_ocr": image_parse_metadata.get("needs_ocr") is True,
        "folder_image_count": sum(item.path.suffix == ".png" for item in folder_plan.items),
        "figure_locator_count": len(figure_assignments),
        "figure_max_chunks_per_id": max(figure_assignments.values(), default=0),
        "large_xlsx_chunk_count": len(large_xlsx_chunks),
        "large_xlsx_heading_only_chunks": sum(
            chunk.text.lstrip().startswith("## Sheet:") and "|" not in chunk.text
            for chunk in large_xlsx_chunks
        ),
        "large_xlsx_table_chunks_without_sheet": sum(
            "|" in chunk.text and "## Sheet:" not in chunk.text
            for chunk in large_xlsx_chunks
        ),
        "large_xlsx_chunks_without_row_locator": sum(
            not chunk.metadata.get("row_range")
            for chunk in large_xlsx_chunks
        ),
        "large_xlsx_duplicate_context_ratio": (
            large_xlsx_repeated_context_tokens / large_xlsx_total_tokens
            if large_xlsx_total_tokens
            else 0.0
        ),
        "large_xlsx_max_token_estimate": max(
            chunk.token_count for chunk in large_xlsx_chunks
        ),
        "large_xlsx_gold_has_row_locator": bool(
            large_gold_chunk.metadata.get("row_range")
        ),
        "large_xlsx_gold_row_contains_300": (
            lambda bounds: bounds[0] <= 300 <= bounds[1]
        )(
            tuple(
                int(value)
                for value in str(large_gold_chunk.metadata["row_range"]).split("-")
            )
        ),
        "large_xlsx_gold_chunk_chars": len(large_gold_chunk.text),
        "dense_chunk_count": len(dense.chunks),
        "dense_max_reference_token_units": max(reference_token_counts),
        "dense_reference_token_budget": token_budget,
        "dense_reference_token_budget_violations": sum(
            count > token_budget for count in reference_token_counts
        ),
        "pdf_two_column_block_order": (
            prepared["pdf-layout"].markdown.index("LEFT ONE")
            < prepared["pdf-layout"].markdown.index("LEFT TWO")
            < prepared["pdf-layout"].markdown.index("RIGHT ONE")
            < prepared["pdf-layout"].markdown.index("RIGHT TWO")
        ),
        "duplicate_content_hash_equal": (
            hashlib.sha256(duplicate).digest() == hashlib.sha256(duplicate).digest()
        ),
        "near_duplicate_content_hash_differs": (
            hashlib.sha256(duplicate).digest()
            != hashlib.sha256(near_duplicate).digest()
        ),
        "same_key_document_identity_stable": same_key_id
        == resolve_document_id(
            namespace="eval",
            collection="docs",
            document_key="shared-key",
            document_id=None,
            policy=DEFAULT_POLICY,
        ),
        "different_key_document_identity_differs": same_key_id != other_key_id,
    }
    metrics.update(await _archive_metrics(tmp_path))
    return metrics


def _assert_document_preparation_contract(
    metrics: dict[str, float | int | bool],
) -> None:
    assert metrics["case_count"] == 14
    assert metrics["exact_fact_retention"] == 1.0
    assert metrics["locator_survival"] >= 0.8
    assert metrics["reliable_span_rate"] == 1.0
    assert metrics["page_specific_ocr_indices"] is True
    assert metrics["mixed_pdf_text_retained"] is True
    assert metrics["image_requires_ocr"] is True
    assert metrics["folder_image_count"] == 0
    assert metrics["figure_locator_count"] == 4
    assert metrics["figure_max_chunks_per_id"] == 1
    assert 2 <= metrics["large_xlsx_chunk_count"] <= 12
    assert metrics["large_xlsx_heading_only_chunks"] == 0
    assert metrics["large_xlsx_table_chunks_without_sheet"] == 0
    assert metrics["large_xlsx_chunks_without_row_locator"] == 0
    assert metrics["large_xlsx_duplicate_context_ratio"] <= 0.10
    assert metrics["large_xlsx_max_token_estimate"] <= 500
    assert metrics["large_xlsx_gold_has_row_locator"] is True
    assert metrics["large_xlsx_gold_row_contains_300"] is True
    assert metrics["dense_reference_token_budget_violations"] == 0
    assert (
        metrics["dense_max_reference_token_units"]
        <= metrics["dense_reference_token_budget"]
    )
    assert metrics["pdf_two_column_block_order"] is True
    assert metrics["archive_item_count"] == 1_000
    assert metrics["archive_memory_item_count"] == 8
    assert metrics["archive_memory_peak_bytes"] < 4_000_000
    assert metrics["duplicate_content_hash_equal"] is True
    assert metrics["near_duplicate_content_hash_differs"] is True
    assert metrics["same_key_document_identity_stable"] is True
    assert metrics["different_key_document_identity_differs"] is True


def test_document_preparation_parser_contract(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("PDF_INSPECTOR_MODE", "disabled")
    metrics = asyncio.run(collect_document_preparation_metrics(tmp_path))
    if os.environ.get("RAG_CORE_PRINT_INGEST_EVAL") == "1":
        print(json.dumps(metrics, indent=2, sort_keys=True))

    _assert_document_preparation_contract(metrics)
