"""Semantic retrieval eval: real local embeddings on a small fixed corpus."""

from __future__ import annotations

import asyncio
import json
import math
import os
import uuid
from collections.abc import Mapping, Sequence
from dataclasses import dataclass, replace
from io import BytesIO
from pathlib import Path
from statistics import mean
from time import perf_counter
from typing import Any, Literal, cast

import pytest

from rag_core import Config
from rag_core.core import Engine
from rag_core.config import (
    LOCAL_EMBEDDING_DIMENSIONS,
    LOCAL_EMBEDDING_MODEL,
    LOCAL_EMBEDDING_PROVIDER,
)
from rag_core.core_models import PreparedDocument
from rag_core.documents.chunking.budget import token_budget_for_char_limit
from rag_core.evals import EvalResult, load_cases, run_eval
from rag_core.events.types import AuditContext
from rag_core.retrieval_defaults import (
    DEFAULT_RERANK,
    DEFAULT_SEARCH_LIMIT,
    DEFAULT_USE_LEXICAL_SEARCH,
)
from rag_core.search import (
    Filter,
    QueryPlan,
    RerankBudget,
    SearchResult,
)
from rag_core.search.context_pack import build_context_pack
from rag_core.search.providers.sparse import FastEmbedSparseEmbedder

pytestmark = [pytest.mark.eval]

_FIXTURE_DIR = Path(__file__).resolve().parent / "semantic_corpus"
_CORPUS_PATH = _FIXTURE_DIR / "corpus.jsonl"
_CASES_PATH = _FIXTURE_DIR / "cases.jsonl"
_NAMESPACE = "semantic_eval"
_COLLECTION = "docs"
_MAX_SECONDS = 120.0
_HEALTHY_RECALL_AT_5_CEILING = 0.97
_HEALTHY_MRR_CEILING = 0.95
_MIN_DEGRADATION_MARGIN = 0.05
_UNRELATED_QUERY = "remote payroll policy for office laptop reimbursement approvals"
_SPARSE_TREATMENT = "dense_plus_explicit_bm25"
_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
    b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
    b"\x00\x00\x0cIDATx\x9cc```\x00\x00\x00\x04\x00\x01"
    b"\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
)

DegradationMode = Literal["rank_shuffle", "wrong_query", "scope_widen"]


@dataclass(frozen=True)
class CorpusDoc:
    document_id: str
    collection: str
    filename: str
    mime_type: str
    file_bytes: bytes


@dataclass(frozen=True)
class AggregateMetrics:
    recall_at_5: float
    recall_at_10: float
    mrr: float
    ndcg_at_10: float


@dataclass(frozen=True)
class MultiSearchMetric:
    name: str
    search_calls: int
    max_search_calls: int
    expected_pass: bool
    forbidden_count: int


@dataclass(frozen=True)
class FormatRetrievalMetric:
    name: str
    rank: int | None
    locator_key: str | None
    locator_value: object
    locator_pass: bool | None


@dataclass(frozen=True)
class _MultiSearchCase:
    name: str
    queries: tuple[str, ...]
    max_search_calls: int
    expected_ids: frozenset[str]
    forbidden_ids: frozenset[str] = frozenset()
    document_ids: tuple[str, ...] = ()


@dataclass(frozen=True)
class SemanticControls:
    sparse_treatment: str
    sparse_diagnostics: Mapping[str, object]
    xlsx_chunk_count: int
    xlsx_duplicate_context_ratio: float
    xlsx_gold_row_range: str | None
    xlsx_max_estimated_tokens: int
    xlsx_max_model_tokens: int
    tokenizer_model: str
    duplicate_raw_exact_hits: int
    duplicate_context_exact_snippets: int
    duplicate_near_context_snippets: int
    redundant_context_chars_avoided: int
    figure_rank: int | None
    figure_id: str | None
    figure_locator_pass: bool
    figure_neighbor_pass: bool
    cross_format: tuple[FormatRetrievalMetric, ...]
    multi_search_workflows: tuple[MultiSearchMetric, ...]
    declared_budget_rejected: bool
    extra_actual_call_rejected: bool


@dataclass(frozen=True)
class SemanticRun:
    results: list[EvalResult]
    controls: SemanticControls


_FLOORS = AggregateMetrics(
    recall_at_5=0.95,
    recall_at_10=0.95,
    mrr=0.83,
    ndcg_at_10=0.85,
)


def test_semantic_corpus_fixture_shape() -> None:
    corpus = _load_corpus(_CORPUS_PATH)
    cases = load_cases(_CASES_PATH)
    document_keys = {doc.document_id for doc in corpus}

    assert 24 <= len(corpus) <= 40
    assert 12 <= len(cases) <= 28
    assert len(document_keys) == len(corpus)
    assert all(case.namespace == _NAMESPACE for case in cases)
    assert all(
        set(case.collections) <= {_COLLECTION, "restricted", "empty"} for case in cases
    )
    assert all(set(case.expected_ids) <= document_keys for case in cases)
    assert any(len(case.expected_ids) > 1 for case in cases)
    content_by_key = {doc.document_id: doc.file_bytes for doc in corpus}
    assert (
        content_by_key["archive_metadata_authority"]
        == content_by_key["archive_metadata_authority_copy"]
    )
    assert (
        content_by_key["archive_metadata_authority"]
        != content_by_key["archive_metadata_authority_variant"]
    )


def test_local_semantic_eval_holds_regression_floors() -> None:
    started = perf_counter()
    run = _run_or_skip()
    elapsed = perf_counter() - started

    assert elapsed <= _MAX_SECONDS, (
        f"local semantic eval took {elapsed:.1f}s; keep the fixed corpus under "
        f"{_MAX_SECONDS:.0f}s post-download"
    )
    _assert_semantic_floors(run.results)
    _assert_healthy_metrics_are_non_trivial(run.results)
    _assert_source_key_scope_and_format_cases(run.results)
    _assert_semantic_controls(run.controls)


def test_local_semantic_eval_rejects_degraded_rankings() -> None:
    for degradation in ("rank_shuffle", "wrong_query"):
        started = perf_counter()
        run = _run_or_skip(degradation=degradation)
        elapsed = perf_counter() - started

        assert elapsed <= _MAX_SECONDS, (
            f"{degradation} semantic eval took {elapsed:.1f}s; keep degradation "
            f"proofs under {_MAX_SECONDS:.0f}s post-download"
        )
        _assert_no_search_errors(run.results)
        failures = _floor_failures(_aggregate(run.results))
        assert failures, (
            f"{degradation} degradation still met semantic floors "
            f"{_format_metrics(_aggregate(run.results))}; harden the corpus or "
            "recalibrate floors before accepting new defaults"
        )
        if degradation == "rank_shuffle":
            _assert_shuffle_margin(_aggregate(run.results))


def test_local_semantic_eval_rejects_scope_widening() -> None:
    run = _run_or_skip(degradation="scope_widen")
    _assert_no_search_errors(run.results)

    controlled = [
        result
        for result in run.results
        if result.case.document_ids
        or result.case.forbidden_ids
        or result.case.forbidden_context_contains
        or result.case.forbidden_private_identifiers
        or result.case.expect_no_results
    ]
    assert controlled
    assert any(
        result.forbidden_id_count > 0
        or result.forbidden_leak_count > 0
        or not result.no_result_pass
        for result in controlled
    )
    content_leak = {result.case.case_id: result for result in run.results}[
        "archive_public_scope_content_safety"
    ]
    assert content_leak.forbidden_id_count == 0
    assert content_leak.forbidden_leak_count > 0
    assert content_leak.prompt_safety_pass is False


def test_local_semantic_eval_skip_fails_in_ci(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("CI", "true")
    monkeypatch.setenv("RAG_CORE_SKIP_FASTEMBED_DOWNLOAD", "1")

    with pytest.raises(BaseException) as exc_info:
        _run_or_skip()

    assert type(exc_info.value).__name__ == "Failed"
    assert "CI must run the local semantic eval gate" in str(exc_info.value)


def _run_or_skip(degradation: DegradationMode | None = None) -> SemanticRun:
    if os.environ.get("RAG_CORE_SKIP_FASTEMBED_DOWNLOAD") == "1":
        _skip_or_fail_unavailable_semantic_gate(
            "RAG_CORE_SKIP_FASTEMBED_DOWNLOAD=1 skips FastEmbed model download"
        )
    try:
        return asyncio.run(_run_semantic_eval(degradation=degradation))
    except Exception as exc:
        if _is_fastembed_download_failure(exc):
            _skip_or_fail_unavailable_semantic_gate(
                "FastEmbed local semantic eval requires one-time model download/cache "
                f"access: {exc}"
            )
        raise


def _skip_or_fail_unavailable_semantic_gate(reason: str) -> None:
    if _truthy_env("CI"):
        pytest.fail(f"CI must run the local semantic eval gate, not skip it: {reason}")
    pytest.skip(reason)


def _truthy_env(name: str) -> bool:
    value = os.environ.get(name, "")
    return value.lower() not in {"", "0", "false", "no", "off"}


async def _run_semantic_eval(
    degradation: DegradationMode | None = None,
) -> SemanticRun:
    config = _semantic_config()
    corpus = _load_corpus(_CORPUS_PATH)
    sparse = FastEmbedSparseEmbedder(enable_splade=False)
    async with Engine(config, sparse_embedder=sparse) as core:
        _assert_local_embedding_runtime(core)
        xlsx_prepared: PreparedDocument | None = None
        content_hashes: dict[str, str] = {}
        for doc in corpus:
            if doc.document_id == "inventory_workbook":
                xlsx_prepared = await core.prepare_bytes(
                    file_bytes=doc.file_bytes,
                    filename=doc.filename,
                    mime_type=doc.mime_type,
                    namespace=_NAMESPACE,
                    collection=doc.collection,
                    document_id=f"semantic-document:{doc.document_id}",
                )
            ingested = await core.add_bytes(
                file_bytes=doc.file_bytes,
                filename=doc.filename,
                mime_type=doc.mime_type,
                namespace=_NAMESPACE,
                collection=doc.collection,
                document_id=f"semantic-document:{doc.document_id}",
                document_key=doc.document_id,
            )
            if ingested.content_sha256 is not None:
                content_hashes[doc.document_id] = ingested.content_sha256
        assert xlsx_prepared is not None
        search_counter = _SearchCallCounter(core)
        eval_core = cast(Engine, search_counter)
        if degradation is not None:
            eval_core = cast(
                Engine,
                _DegradedSearchCore(eval_core, degradation),
            )
        results = await run_eval(
            eval_core,
            load_cases(_CASES_PATH),
            k_values=(5, 10, len(corpus)),
            rerank=False,
        )
        controls = await _collect_semantic_controls(
            eval_core=eval_core,
            xlsx_prepared=xlsx_prepared,
            content_hashes=content_hashes,
            sparse=sparse,
            search_counter=search_counter,
        )
        return SemanticRun(results=results, controls=controls)


class _SearchCallCounter:
    def __init__(self, core: Engine) -> None:
        self._core = core
        self.call_count = 0

    async def search(self, **kwargs: object) -> list[SearchResult]:
        self.call_count += 1
        return cast(
            list[SearchResult],
            await cast(Any, self._core).search(**kwargs),
        )


class _DuplicateSearchCore:
    def __init__(self, core: Engine) -> None:
        self._core = core

    async def search(self, **kwargs: object) -> list[SearchResult]:
        first = await cast(Any, self._core).search(**kwargs)
        await cast(Any, self._core).search(**kwargs)
        return cast(list[SearchResult], first)


class _DegradedSearchCore:
    def __init__(self, core: Engine, degradation: DegradationMode) -> None:
        self._core = core
        self._degradation = degradation

    async def search(
        self,
        *,
        query: str,
        namespace: str,
        collections: list[str],
        limit: int = DEFAULT_SEARCH_LIMIT,
        content_types: list[str] | None = None,
        document_ids: list[str] | None = None,
        rerank: bool = DEFAULT_RERANK,
        use_lexical_search: bool = DEFAULT_USE_LEXICAL_SEARCH,
        query_plan: QueryPlan | None = None,
        metadata_filter: Filter | None = None,
        rerank_budget: RerankBudget | None = None,
        audit_context: AuditContext | None = None,
    ) -> list[SearchResult]:
        search_query = _UNRELATED_QUERY if self._degradation == "wrong_query" else query
        resolved_collections = (
            ["docs", "restricted"]
            if self._degradation == "scope_widen"
            else collections
        )
        resolved_document_ids = (
            None if self._degradation == "scope_widen" else document_ids
        )
        hits = await self._core.search(
            query=search_query,
            namespace=namespace,
            collections=resolved_collections,
            limit=limit,
            content_types=content_types,
            document_ids=resolved_document_ids,
            rerank=rerank,
            use_lexical_search=use_lexical_search,
            query_plan=query_plan,
            metadata_filter=metadata_filter,
            rerank_budget=rerank_budget,
            audit_context=audit_context,
        )
        if self._degradation == "rank_shuffle":
            return list(reversed(hits))
        return hits


def _semantic_config() -> Config:
    config = Config.local()
    return replace(
        config,
        vector_store=replace(
            config.vector_store,
            qdrant=replace(
                config.vector_store.qdrant,
                store_collection=f"rag_core_semantic_eval_{uuid.uuid4().hex}",
            ),
        ),
    )


def _assert_local_embedding_runtime(core: Engine) -> None:
    embedding = core.describe_runtime()["embedding"]
    assert embedding == {
        "provider": LOCAL_EMBEDDING_PROVIDER,
        "model": LOCAL_EMBEDDING_MODEL,
        "dimensions": LOCAL_EMBEDDING_DIMENSIONS,
    }, (
        "Local semantic eval model changed; rerun 5x calibration, degradation "
        f"demos, and update floors before accepting new defaults. got={embedding!r}"
    )


def _load_corpus(path: Path) -> list[CorpusDoc]:
    docs: list[CorpusDoc] = []
    with path.open(encoding="utf-8") as handle:
        for raw in handle:
            row = json.loads(raw)
            document_id = str(row["document_id"])
            fixture = row.get("fixture")
            if fixture == "large_xlsx":
                filename = f"{document_id}.xlsx"
                mime_type = (
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
                file_bytes = _large_xlsx_bytes()
            elif fixture == "two_figure_docx":
                filename = f"{document_id}.docx"
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                file_bytes = _figure_docx_bytes()
            elif fixture == "locator_pdf":
                filename = f"{document_id}.pdf"
                mime_type = "application/pdf"
                file_bytes = _locator_pdf_bytes()
            elif fixture == "locator_pptx":
                filename = f"{document_id}.pptx"
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                file_bytes = _locator_pptx_bytes()
            elif fixture == "structured_jsonl":
                filename = f"{document_id}.jsonl"
                mime_type = "application/x-ndjson"
                file_bytes = _structured_jsonl_bytes()
            else:
                filename = f"{document_id}.md"
                mime_type = "text/markdown"
                file_bytes = f"# {row['title']}\n\n{row['body']}".encode()
            docs.append(
                CorpusDoc(
                    document_id=document_id,
                    collection=str(row.get("collection", _COLLECTION)),
                    filename=filename,
                    mime_type=mime_type,
                    file_bytes=file_bytes,
                )
            )
    return docs


def _large_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Supplier Inventory"
    sheet.append(["record", "state", "inspection"])
    for index in range(1, 301):
        inspection = (
            "supplier lot KQ-733 is quarantined pending chromatic inspection"
            if index == 299
            else f"routine receiving check {index:03d}"
        )
        sheet.append([f"lot-{index:03d}", "received", inspection])
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _figure_docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("Acetate Flattening Workflow", level=1)
    target = document.add_paragraph()
    target.add_run(
        "ACETATE-BEFORE Technicians humidify curled acetate in a sealed chamber "
        "and flatten it beneath inert polyester. "
    )
    target_shape = target.add_run().add_picture(BytesIO(_PNG))
    target_shape._inline.docPr.attrib["descr"] = "Acetate flattening workflow diagram"
    target.add_run(
        " ACETATE-AFTER The flattened sheet moves to an overhead camera with "
        "diffuse copy lights."
    )
    document.add_paragraph("Unrelated register-description separator material. " * 100)
    decoy = document.add_paragraph()
    decoy.add_run(
        "MAP-SLEEVE-BEFORE Staff inspect buffered sleeves for oversize maps. "
    )
    decoy_shape = decoy.add_run().add_picture(BytesIO(_PNG))
    decoy_shape._inline.docPr.attrib["descr"] = "Map sleeve inspection diagram"
    decoy.add_run(" MAP-SLEEVE-AFTER The sleeve record is filed with the housing log.")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _locator_pdf_bytes() -> bytes:
    import fitz

    document = fitz.open()
    first = document.new_page()
    first.insert_text((72, 60), "GENERAL COLD STORAGE OVERVIEW")
    for index in range(8):
        first.insert_text(
            (72, 100 + index * 24),
            f"Routine cabinet orientation record {index + 1}.",
        )
    second = document.new_page()
    second.insert_text((72, 60), "Nitrate Negative Transfer")
    for index in range(8):
        second.insert_text(
            (72, 100 + index * 24),
            (
                "After room-temperature acclimation, nitrate negatives move "
                f"into a vented cold cabinet for stable preservation {index + 1}."
            ),
        )
    payload = cast(bytes, document.tobytes())
    document.close()
    return payload


def _locator_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    if first.shapes.title is not None:
        first.shapes.title.text = "General Exhibit Review"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    if second.shapes.title is not None:
        second.shapes.title.text = "Fragile Mount Decision"
    box = second.shapes.add_textbox(
        Inches(1),
        Inches(1.5),
        Inches(8),
        Inches(2),
    )
    box.text_frame.text = (
        "Use a low-angle cradle with magnetic edge strips when mounting the "
        "fragile display sheet."
    )
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _structured_jsonl_bytes() -> bytes:
    return (
        b'{"sensor":"north-vault","event":"humidity check","status":"normal"}\n'
        b'{"sensor":"east-shelf","event":"vibration alert",'
        b'"resolution":"isolated shelf fan"}\n'
    )


async def _collect_semantic_controls(
    *,
    eval_core: Engine,
    xlsx_prepared: PreparedDocument,
    content_hashes: Mapping[str, str],
    sparse: FastEmbedSparseEmbedder,
    search_counter: _SearchCallCounter,
) -> SemanticControls:
    xlsx_hits = await eval_core.search(
        query=("late inventory exception requiring supplier chromatic inspection"),
        namespace=_NAMESPACE,
        collections=[_COLLECTION],
        limit=10,
        rerank=False,
    )
    xlsx_hit = next(
        (hit for hit in xlsx_hits if hit.document_key == "inventory_workbook"),
        None,
    )
    xlsx_gold_row_range = (
        str(xlsx_hit.metadata.get("row_range"))
        if xlsx_hit is not None and xlsx_hit.metadata.get("row_range") is not None
        else None
    )
    total_xlsx_tokens = sum(chunk.token_count for chunk in xlsx_prepared.chunks)
    repeated_context_tokens = sum(
        _metadata_int(chunk.metadata.get("xlsx_context_token_count"))
        for chunk in xlsx_prepared.chunks[1:]
    )
    model_token_counts = await asyncio.to_thread(
        _independent_model_token_counts,
        [chunk.embedding_text for chunk in xlsx_prepared.chunks],
    )

    duplicate_hits = await eval_core.search(
        query=(
            "where alternate creator spellings become one preferred catalog identity"
        ),
        namespace=_NAMESPACE,
        collections=[_COLLECTION],
        limit=10,
        rerank=False,
    )
    duplicate_pack = build_context_pack(
        duplicate_hits,
        query="preferred catalog creator identity",
        max_snippets=10,
    )
    exact_hash = content_hashes["archive_metadata_authority"]
    near_hash = content_hashes["archive_metadata_authority_variant"]
    exact_raw = [hit for hit in duplicate_hits if hit.content_sha256 == exact_hash]
    exact_context = sum(
        snippet.source.content_sha256 == exact_hash
        for snippet in duplicate_pack.snippets
    )
    near_context = sum(
        snippet.source.content_sha256 == near_hash
        for snippet in duplicate_pack.snippets
    )
    redundant_chars = sum(len(hit.text) for hit in exact_raw) - max(
        (len(hit.text) for hit in exact_raw), default=0
    )

    figure_hits = await eval_core.search(
        query=(
            "illustrated process for flattening curled acetate before overhead imaging"
        ),
        namespace=_NAMESPACE,
        collections=[_COLLECTION],
        limit=10,
        rerank=False,
    )
    figure_entry = next(
        (
            (rank, hit)
            for rank, hit in enumerate(figure_hits, start=1)
            if hit.document_key == "archive_figure_workflow"
            and (
                hit.figure_id == "fig:docx:1"
                or hit.metadata.get("figure_id") == "fig:docx:1"
            )
        ),
        None,
    )

    cross_format = (
        await _format_retrieval_metric(
            eval_core,
            name="pdf_section_locator",
            query=(
                "file describing nitrate negatives entering ventilated cold "
                "storage after acclimation"
            ),
            document_key="archive_cold_storage_pdf",
            locator_key="section_title",
            expected_locator="Nitrate Negative Transfer",
        ),
        await _format_retrieval_metric(
            eval_core,
            name="pptx_slide_locator",
            query=(
                "presentation about a fragile display supported by a low-angle "
                "cradle and magnetic strips"
            ),
            document_key="archive_mount_review_pptx",
            locator_key="slide_number",
            expected_locator=2,
        ),
        await _format_retrieval_metric(
            eval_core,
            name="jsonl_rank_only_retrieval",
            query=(
                "structured preservation event where an east shelf vibration "
                "alert ended after fan isolation"
            ),
            document_key="archive_sensor_events_jsonl",
            locator_key=None,
            expected_locator=None,
        ),
    )

    multi_search_cases = (
        _MultiSearchCase(
            name="transit_inspection_and_release",
            queries=(
                "record that quantifies stopping performance on the service track",
                "record containing the signoff that returns a tram to passengers",
            ),
            max_search_calls=2,
            expected_ids=frozenset(
                {"transit_brake_inspection", "transit_service_release_log"}
            ),
            forbidden_ids=frozenset({"archive_emergency_lighting"}),
        ),
        _MultiSearchCase(
            name="document_scoped_brake_search",
            queries=("where is stopping distance measured before release",),
            max_search_calls=1,
            expected_ids=frozenset({"transit_brake_inspection"}),
            forbidden_ids=frozenset(
                {"transit_service_release_log", "archive_emergency_lighting"}
            ),
            document_ids=("semantic-document:transit_brake_inspection",),
        ),
    )
    multi_search_metrics: list[MultiSearchMetric] = []
    for case in multi_search_cases:
        multi_search_metrics.append(
            await _run_bounded_multi_search(
                eval_core,
                case,
                search_counter=search_counter,
            )
        )
    try:
        await _run_bounded_multi_search(
            eval_core,
            _MultiSearchCase(
                name="over_budget_control",
                queries=("one", "two", "three"),
                max_search_calls=2,
                expected_ids=frozenset(),
            ),
            search_counter=search_counter,
        )
        declared_budget_rejected = False
    except ValueError:
        declared_budget_rejected = True

    try:
        await _run_bounded_multi_search(
            cast(Engine, _DuplicateSearchCore(eval_core)),
            _MultiSearchCase(
                name="extra_actual_call_control",
                queries=("record with measured stopping performance",),
                max_search_calls=1,
                expected_ids=frozenset({"transit_brake_inspection"}),
            ),
            search_counter=search_counter,
        )
        extra_actual_call_rejected = False
    except ValueError:
        extra_actual_call_rejected = True

    return SemanticControls(
        sparse_treatment=_SPARSE_TREATMENT,
        sparse_diagnostics=sparse.diagnostics(),
        xlsx_chunk_count=len(xlsx_prepared.chunks),
        xlsx_duplicate_context_ratio=(
            repeated_context_tokens / total_xlsx_tokens if total_xlsx_tokens else 0.0
        ),
        xlsx_gold_row_range=xlsx_gold_row_range,
        xlsx_max_estimated_tokens=max(
            (chunk.token_count for chunk in xlsx_prepared.chunks),
            default=0,
        ),
        xlsx_max_model_tokens=max(model_token_counts, default=0),
        tokenizer_model=LOCAL_EMBEDDING_MODEL,
        duplicate_raw_exact_hits=len(exact_raw),
        duplicate_context_exact_snippets=exact_context,
        duplicate_near_context_snippets=near_context,
        redundant_context_chars_avoided=max(0, redundant_chars),
        figure_rank=figure_entry[0] if figure_entry is not None else None,
        figure_id=(figure_entry[1].figure_id if figure_entry is not None else None),
        figure_locator_pass=(
            figure_entry is not None and figure_entry[1].figure_id == "fig:docx:1"
        ),
        figure_neighbor_pass=(
            figure_entry is not None
            and "ACETATE-BEFORE" in figure_entry[1].text
            and "ACETATE-AFTER" in figure_entry[1].text
            and "MAP-SLEEVE-BEFORE" not in figure_entry[1].text
        ),
        cross_format=cross_format,
        multi_search_workflows=tuple(multi_search_metrics),
        declared_budget_rejected=declared_budget_rejected,
        extra_actual_call_rejected=extra_actual_call_rejected,
    )


async def _format_retrieval_metric(
    core: Engine,
    *,
    name: str,
    query: str,
    document_key: str,
    locator_key: str | None,
    expected_locator: object,
) -> FormatRetrievalMetric:
    hits = await core.search(
        query=query,
        namespace=_NAMESPACE,
        collections=[_COLLECTION],
        limit=10,
        rerank=False,
    )
    entry = next(
        (
            (rank, hit)
            for rank, hit in enumerate(hits, start=1)
            if hit.document_key == document_key
            and (
                locator_key is None
                or (getattr(hit, locator_key, None) or hit.metadata.get(locator_key))
                == expected_locator
            )
        ),
        None,
    )
    locator_value = (
        (getattr(entry[1], locator_key, None) or entry[1].metadata.get(locator_key))
        if entry is not None and locator_key is not None
        else None
    )
    return FormatRetrievalMetric(
        name=name,
        rank=entry[0] if entry is not None else None,
        locator_key=locator_key,
        locator_value=locator_value,
        locator_pass=(entry is not None if locator_key is not None else None),
    )


async def _run_bounded_multi_search(
    core: Engine,
    case: _MultiSearchCase,
    *,
    search_counter: _SearchCallCounter,
) -> MultiSearchMetric:
    starting_calls = search_counter.call_count
    found_ids: set[str] = set()
    forbidden_count = 0
    for query in case.queries:
        if search_counter.call_count - starting_calls >= case.max_search_calls:
            raise ValueError("bounded multi-search call budget exceeded")
        hits = await core.search(
            query=query,
            namespace=_NAMESPACE,
            collections=[_COLLECTION],
            document_ids=list(case.document_ids) or None,
            limit=5,
            rerank=False,
        )
        if search_counter.call_count - starting_calls > case.max_search_calls:
            raise ValueError("bounded multi-search call budget exceeded")
        hit_ids = {hit.document_key or hit.document_id or hit.id for hit in hits}
        found_ids.update(hit_ids)
        forbidden_count += len(hit_ids & case.forbidden_ids)
    return MultiSearchMetric(
        name=case.name,
        search_calls=search_counter.call_count - starting_calls,
        max_search_calls=case.max_search_calls,
        expected_pass=case.expected_ids <= found_ids,
        forbidden_count=forbidden_count,
    )


def _independent_model_token_counts(texts: list[str]) -> list[int]:
    from fastembed import TextEmbedding

    embedding = TextEmbedding(
        model_name=LOCAL_EMBEDDING_MODEL,
        local_files_only=True,
    )
    tokenizer = getattr(getattr(embedding, "model"), "tokenizer")
    getattr(tokenizer, "no_truncation")()
    return [len(getattr(tokenizer.encode(text), "ids")) for text in texts]


def _metadata_int(value: object) -> int:
    return value if isinstance(value, int) and not isinstance(value, bool) else 0


def _aggregate(results: Sequence[EvalResult]) -> AggregateMetrics:
    answerable = [result for result in results if result.case.expected_ids]
    assert answerable
    return AggregateMetrics(
        recall_at_5=mean(result.recall_at_5 for result in answerable),
        recall_at_10=mean(result.recall_at_10 for result in answerable),
        mrr=mean(result.mrr for result in answerable),
        ndcg_at_10=mean(result.ndcg_at_10 for result in answerable),
    )


def _assert_semantic_floors(
    results: list[EvalResult],
) -> None:
    _assert_no_search_errors(results)

    metrics = _aggregate(results)
    for (metric_name, actual), (_, floor) in zip(
        _metric_pairs(metrics),
        _metric_pairs(_FLOORS),
        strict=True,
    ):
        assert actual >= floor, _floor_message(metric_name, actual, floor)


def _assert_no_search_errors(results: Sequence[EvalResult]) -> None:
    errors = {
        result.case.case_id or result.case.query: result.error_type
        for result in results
        if result.error_type is not None
    }
    assert not errors, f"semantic eval search errors: {errors}"


def _assert_healthy_metrics_are_non_trivial(results: Sequence[EvalResult]) -> None:
    metrics = _aggregate(results)
    imperfect_count = _imperfect_case_count(results)
    answerable_count = sum(bool(result.case.expected_ids) for result in results)
    required = math.ceil(answerable_count / 5)
    assert imperfect_count >= required, (
        f"semantic corpus has only {imperfect_count}/{answerable_count} "
        "imperfect answerable cases; "
        "harden near-miss cases and rerun 5x calibration before updating floors"
    )
    assert (
        metrics.recall_at_5 <= _HEALTHY_RECALL_AT_5_CEILING
        or metrics.mrr <= _HEALTHY_MRR_CEILING
    ), (
        "semantic corpus became trivially separable "
        f"{_format_metrics(metrics)}; harden near-miss cases and rerun 5x "
        "calibration before updating floors"
    )


def _assert_source_key_scope_and_format_cases(
    results: Sequence[EvalResult],
) -> None:
    by_case = {result.case.case_id: result for result in results}
    duplicate_case = by_case["archive_preferred_creator_names"]
    assert duplicate_case.mrr >= 0.5

    scoped_case = by_case["archive_restricted_emergency_lighting"]
    assert scoped_case.retrieved_ids[0] == "archive_emergency_lighting"
    assert scoped_case.mrr == 1.0

    xlsx_case = by_case["inventory_late_supplier_exception"]
    assert xlsx_case.recall_at_5 == 1.0
    assert xlsx_case.mrr >= 0.5

    positive_scope = by_case["transit_brake_document_scope"]
    assert positive_scope.retrieved_ids == ("transit_brake_inspection",)
    assert positive_scope.forbidden_id_count == 0

    negative_scope = by_case["transit_brake_negative_document_scope"]
    assert negative_scope.no_result_pass is True
    assert negative_scope.forbidden_id_count == 0

    empty_scope = by_case["empty_collection_scope_control"]
    assert empty_scope.no_result_pass is True
    assert empty_scope.retrieved_ids == ()
    # Populated-corpus abstention needs a score or answer contract that Engine
    # does not currently own; this case only proves an explicitly empty scope.

    figure_case = by_case["archive_inline_figure_retrieval"]
    assert figure_case.recall_at_5 == 1.0
    public_scope = by_case["archive_public_scope_content_safety"]
    assert public_scope.forbidden_leak_count == 0
    assert public_scope.prompt_safety_pass is True
    for case_id in (
        "pdf_cold_storage_locator",
        "pptx_exhibit_mount_locator",
        "jsonl_sensor_event_retrieval",
    ):
        assert by_case[case_id].recall_at_5 == 1.0
    assert all(result.forbidden_id_count == 0 for result in results)


def _assert_semantic_controls(controls: SemanticControls) -> None:
    assert controls.sparse_treatment == _SPARSE_TREATMENT
    assert controls.sparse_diagnostics == {
        "provider": "fastembed",
        "bm25_enabled": True,
        "bm25_load_status": "loaded",
        "splade_enabled": False,
        "splade_load_status": "disabled",
    }

    assert 2 <= controls.xlsx_chunk_count <= 12
    assert controls.xlsx_duplicate_context_ratio <= 0.10
    assert controls.xlsx_gold_row_range is not None
    start_row, end_row = (
        int(value) for value in controls.xlsx_gold_row_range.split("-")
    )
    assert start_row <= 300 <= end_row
    assert controls.xlsx_max_estimated_tokens <= token_budget_for_char_limit(2_000)
    assert controls.tokenizer_model == LOCAL_EMBEDDING_MODEL
    assert controls.xlsx_max_model_tokens <= 512

    assert controls.duplicate_raw_exact_hits >= 2
    assert controls.duplicate_context_exact_snippets == 1
    assert controls.duplicate_near_context_snippets >= 1
    assert controls.redundant_context_chars_avoided > 0

    assert controls.figure_rank is not None
    assert controls.figure_rank <= 5
    assert controls.figure_id == "fig:docx:1"
    assert controls.figure_locator_pass is True
    assert controls.figure_neighbor_pass is True
    assert {metric.name for metric in controls.cross_format} == {
        "pdf_section_locator",
        "pptx_slide_locator",
        "jsonl_rank_only_retrieval",
    }
    assert all(
        metric.rank is not None and metric.rank <= 5 for metric in controls.cross_format
    )
    locator_metrics = [
        metric for metric in controls.cross_format if metric.locator_key is not None
    ]
    assert {metric.name for metric in locator_metrics} == {
        "pdf_section_locator",
        "pptx_slide_locator",
    }
    assert all(metric.locator_pass is True for metric in locator_metrics)
    [jsonl_metric] = [
        metric
        for metric in controls.cross_format
        if metric.name == "jsonl_rank_only_retrieval"
    ]
    assert jsonl_metric.locator_pass is None
    assert controls.declared_budget_rejected is True
    assert controls.extra_actual_call_rejected is True
    assert controls.multi_search_workflows
    assert all(
        metric.search_calls <= metric.max_search_calls
        and metric.expected_pass
        and metric.forbidden_count == 0
        for metric in controls.multi_search_workflows
    )


def _imperfect_case_count(results: Sequence[EvalResult]) -> int:
    return sum(
        1
        for result in results
        if result.case.expected_ids
        if (result.recall_at_5 < 1.0 or result.mrr < 1.0 or result.ndcg_at_10 < 1.0)
    )


def _floor_failures(metrics: AggregateMetrics) -> tuple[str, ...]:
    failures: list[str] = []
    for (metric_name, actual), (_, floor) in zip(
        _metric_pairs(metrics),
        _metric_pairs(_FLOORS),
        strict=True,
    ):
        if actual < floor:
            failures.append(f"{metric_name}={actual:.3f}<floor={floor:.3f}")
    return tuple(failures)


def _assert_shuffle_margin(metrics: AggregateMetrics) -> None:
    for (metric_name, degraded), (_, floor) in zip(
        _metric_pairs(metrics),
        _metric_pairs(_FLOORS),
        strict=True,
    ):
        assert floor >= degraded + _MIN_DEGRADATION_MARGIN, (
            f"semantic eval {metric_name} floor={floor:.3f} must stay at least "
            f"{_MIN_DEGRADATION_MARGIN:.2f} above rank-shuffle degradation "
            f"{degraded:.3f}; harden the corpus or recalibrate floors"
        )


def _metric_pairs(metrics: AggregateMetrics) -> tuple[tuple[str, float], ...]:
    return (
        ("recall@5", metrics.recall_at_5),
        ("recall@10", metrics.recall_at_10),
        ("mrr", metrics.mrr),
        ("ndcg@10", metrics.ndcg_at_10),
    )


def _format_metrics(metrics: AggregateMetrics) -> str:
    return " ".join(
        f"{metric_name}={value:.3f}" for metric_name, value in _metric_pairs(metrics)
    )


def _floor_message(metric: str, actual: float, floor: float) -> str:
    return (
        f"semantic eval {metric}={actual:.3f} below floor={floor:.3f}; "
        "inspect retrieval changes, then rerun 5x calibration and degradation demos "
        "before updating floors"
    )


def _is_fastembed_download_failure(exc: Exception) -> bool:
    text = str(exc).lower()
    return (
        "fastembed local embedding model failed to load" in text
        or "huggingface" in text
        or "model download" in text
        or "could not download" in text
    )
